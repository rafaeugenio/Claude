"""
template_base.py — Gerador de apresentações PowerPoint Citiesoft
Baseia-se no arquivo Onboarding.pptx como template visual.

PRÉ-REQUISITO:
  1. Copie assets/Onboarding.pptx para /home/claude/Onboarding.pptx
  2. pip install python-pptx

USO:
  Edite apenas o bloco "CONTEÚDO DA APRESENTAÇÃO" no final do arquivo.
  Execute: python3 template_base.py
"""

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from copy import deepcopy

# ─── CONFIGURAÇÃO ─────────────────────────────────────────────────────────────
TEMPLATE_PATH = "/home/claude/Onboarding.pptx"
OUTPUT_PATH   = "/home/claude/citiesoft_apresentacao.pptx"

# ─── CORES OFICIAIS CITIESOFT ─────────────────────────────────────────────────
BLUE      = RGBColor(0x14, 0x6B, 0xFA)  # #146BFA — azul primário da marca
GRAY_DARK = RGBColor(0x3C, 0x3C, 0x3C)  # #3C3C3C — texto principal (fundo branco + logo)
BLACK     = RGBColor(0x00, 0x00, 0x00)  # #000000 — texto em slides de fundo alternativo

# ─── FONTES ───────────────────────────────────────────────────────────────────
FONT_DISPLAY = "Red Hat Display"  # títulos, citações, destaques, bullets-chave
FONT_BODY    = "Inter"            # corpo de texto descritivo
FONT_MONO    = "Red Hat Mono"     # rodapé de seção (pequeno, 5pt)

# ─── ÍNDICES DOS SLIDES TEMPLATE NO Onboarding.pptx (0-indexado) ─────────────
TMPL_QUOTE    = 2   # Slide 3  → Citação / Frase de Impacto
TMPL_CONTENT  = 4   # Slide 5  → Título + Corpo + Imagem lateral
TMPL_KEYWORDS = 7   # Slide 8  → Palavras-chave em azul
TMPL_BULLETS  = 11  # Slide 12 → Bullets / Lista de ações


# ─── HELPERS INTERNOS ─────────────────────────────────────────────────────────

def _clone_slide(prs, slide_idx):
    """
    Clona um slide dentro da mesma apresentação, preservando imagens e logo.
    O clone é adicionado ao final da lista de slides.
    """
    src = prs.slides[slide_idx]
    blank = prs.slide_layouts[6]
    dst = prs.slides.add_slide(blank)

    # Substituir spTree pelo conteúdo do slide fonte
    src_spTree = src.shapes._spTree
    dst_spTree = dst.shapes._spTree
    for child in list(dst_spTree):
        dst_spTree.remove(child)
    for child in src_spTree:
        dst_spTree.append(deepcopy(child))

    # Compartilhar relacionamentos de imagem (rIds) para que as imagens resolvam
    for rId, rel in src.part._rels.items():
        if rId not in dst.part._rels:
            try:
                dst.part._rels[rId] = rel
            except Exception:
                pass

    return dst


def _delete_slide(prs, slide_idx):
    """Remove um slide pelo índice (0-indexado)."""
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[slide_idx])


def _find_textbox_near(slide, left_in, top_in, tol=0.18):
    """Retorna o shape com text frame cuja posição mais se aproxima de (left_in, top_in)."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        l = shape.left / 914400
        t = shape.top / 914400
        if abs(l - left_in) < tol and abs(t - top_in) < tol:
            return shape
    return None


def _set_simple(shape, text, font_name, size_pt, bold, color, align=PP_ALIGN.LEFT):
    """Substitui todo o conteúdo de um text frame por uma run única."""
    tf = shape.text_frame
    # Manter apenas o primeiro parágrafo
    while len(tf.paragraphs) > 1:
        tf.paragraphs[-1]._p.getparent().remove(tf.paragraphs[-1]._p)
    p = tf.paragraphs[0]
    p.alignment = align
    # Remover runs existentes
    for r in list(p._p.findall(qn('a:r'))):
        p._p.remove(r)
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color


def _set_spans(shape, spans, size_pt, align=PP_ALIGN.LEFT):
    """
    Substitui o conteúdo de um text frame por múltiplas runs inline coloridas.
    spans: list of (texto, RGBColor, bold)
    Exemplo: [("Texto comum ", GRAY_DARK, False), ("destaque", BLUE, True), (".", GRAY_DARK, False)]
    """
    tf = shape.text_frame
    while len(tf.paragraphs) > 1:
        tf.paragraphs[-1]._p.getparent().remove(tf.paragraphs[-1]._p)
    p = tf.paragraphs[0]
    p.alignment = align
    for r in list(p._p.findall(qn('a:r'))):
        p._p.remove(r)
    for text, color, bold in spans:
        run = p.add_run()
        run.text = text
        run.font.name = FONT_DISPLAY
        run.font.size = Pt(size_pt)
        run.font.bold = bold
        run.font.color.rgb = color


def _set_multiline(shape, lines, font_name, size_pt, bold, color, align=PP_ALIGN.LEFT):
    """Substitui o conteúdo de um text frame por múltiplos parágrafos (um por item de lines)."""
    tf = shape.text_frame
    while len(tf.paragraphs) > 1:
        tf.paragraphs[-1]._p.getparent().remove(tf.paragraphs[-1]._p)
    p0 = tf.paragraphs[0]
    p0.alignment = align
    for r in list(p0._p.findall(qn('a:r'))):
        p0._p.remove(r)

    for i, line in enumerate(lines):
        p = p0 if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(size_pt)
        run.font.bold = bold
        run.font.color.rgb = color


# ─── TIPOS DE SLIDE ──────────────────────────────────────────────────────────

def quote_slide(prs, spans, section_name):
    """
    Slide de Citação / Frase de Impacto.
    Ideal para frases de impacto, posicionamentos e valores da empresa.

    spans: list of (texto, cor, bold)
      Partes do texto montado inline — misture GRAY_DARK e BLUE para criar destaques.
      Exemplo:
        [
          ('"Onde tecnologia gera ', GRAY_DARK, False),
          ('impacto real', BLUE, True),
          ('" — Citiesoft', GRAY_DARK, False),
        ]

    section_name: str — nome da seção exibido no rodapé esquerdo inferior.
      Exemplo: "Boas vindas à Citiesoft"
    """
    slide = _clone_slide(prs, TMPL_QUOTE)

    # Texto principal da citação (posição original: top ≈ 2.031in)
    main_box = _find_textbox_near(slide, 0.833, 2.031)
    if main_box:
        _set_spans(main_box, spans, size_pt=21.4)

    # Rodapé de seção (posição original: top ≈ 4.719in)
    footer = _find_textbox_near(slide, 0.833, 4.719)
    if footer:
        _set_simple(footer, section_name, FONT_MONO, 5.0, True, GRAY_DARK)

    return slide


def content_slide(prs, title_line1, title_line2, body_text, section_name):
    """
    Slide de Título + Corpo + Imagem lateral.
    Ideal para apresentar um conceito, definição ou explicação com imagem à direita.

    title_line1: str — primeira linha do título (Red Hat Display, preto, sem bold).
      Exemplo: "O que é a"
    title_line2: str — segunda linha do título (Red Hat Display, azul, bold).
      Exemplo: "Citiesoft?"
    body_text: str OU list of str — texto descritivo em Inter.
      Se for list, cada item é um parágrafo separado.
      Exemplo: "A Citiesoft é uma empresa de tecnologia especializada em soluções digitais..."
    section_name: str — nome da seção exibido no rodapé.
    """
    slide = _clone_slide(prs, TMPL_CONTENT)

    # Título (posição original: top ≈ 1.418in) — 2 parágrafos
    title_box = _find_textbox_near(slide, 0.833, 1.418)
    if title_box:
        tf = title_box.text_frame
        while len(tf.paragraphs) > 1:
            tf.paragraphs[-1]._p.getparent().remove(tf.paragraphs[-1]._p)
        p1 = tf.paragraphs[0]
        for r in list(p1._p.findall(qn('a:r'))):
            p1._p.remove(r)
        r1 = p1.add_run()
        r1.text = title_line1
        r1.font.name = FONT_DISPLAY
        r1.font.size = Pt(25.6)
        r1.font.bold = False
        r1.font.color.rgb = BLACK
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = title_line2
        r2.font.name = FONT_DISPLAY
        r2.font.size = Pt(25.6)
        r2.font.bold = True
        r2.font.color.rgb = BLUE

    # Corpo de texto (posição original: top ≈ 2.812in)
    body_box = _find_textbox_near(slide, 0.833, 2.812)
    if body_box:
        if isinstance(body_text, list):
            _set_multiline(body_box, body_text, FONT_BODY, 5.7, False, BLACK)
        else:
            _set_simple(body_box, body_text, FONT_BODY, 5.7, False, BLACK)

    # Rodapé (posição original: top ≈ 4.719in)
    footer = _find_textbox_near(slide, 0.833, 4.719)
    if footer:
        _set_simple(footer, section_name, FONT_MONO, 5.0, True, BLACK)

    return slide


def keywords_slide(prs, page_title, subtitle_label, keywords, body_text, section_name):
    """
    Slide de Palavras-chave em azul.
    Ideal para listar atributos, características, valores ou pilares — palavras curtas
    que impactam visualmente por tamanho e cor.

    page_title: str — título do slide (esquerda, topo).
      Exemplo: "Tom de voz da Citiesoft"
    subtitle_label: str — label descritivo antes das keywords.
      Exemplo: "Como nos comunicamos:"
    keywords: list of str — palavras curtas em azul, max ~6 itens.
      Exemplo: ["Clara", "Confiável", "Institucional", "Tecnológica", "Parceira"]
    body_text: str — texto de apoio em Inter (esquerda, abaixo do label).
    section_name: str — rodapé de seção.
    """
    slide = _clone_slide(prs, TMPL_KEYWORDS)

    # Título da página (posição original: top ≈ 0.833in)
    title_box = _find_textbox_near(slide, 0.833, 0.833)
    if title_box:
        _set_simple(title_box, page_title, FONT_DISPLAY, 25.6, False, BLACK)

    # Label do subtítulo (posição original: top ≈ 2.151in)
    sublabel_box = _find_textbox_near(slide, 0.833, 2.151)
    if sublabel_box:
        _set_simple(sublabel_box, subtitle_label, FONT_DISPLAY, 12.8, False, BLACK)

    # Keywords em azul (posição original: left ≈ 5.000in, top ≈ 2.130in)
    kw_box = _find_textbox_near(slide, 5.000, 2.130)
    if kw_box:
        _set_multiline(kw_box, keywords, FONT_DISPLAY, 25.6, False, BLUE)

    # Texto de apoio (posição original: top ≈ 2.812in)
    body_box = _find_textbox_near(slide, 0.833, 2.812)
    if body_box:
        _set_simple(body_box, body_text, FONT_BODY, 5.7, False, BLACK)

    # Rodapé (posição original: top ≈ 4.719in)
    footer = _find_textbox_near(slide, 0.833, 4.719)
    if footer:
        _set_simple(footer, section_name, FONT_MONO, 5.0, True, BLACK)

    return slide


def bullets_slide(prs, intro_spans, bullets, section_name):
    """
    Slide de Bullets / Lista de Ações.
    Ideal para next steps, checklists, procedimentos ou etapas numeradas.

    intro_spans: list of (texto, cor, bold) — frase introdutória com destaques inline.
      Exemplo:
        [
          ("A partir de agora os ", GRAY_DARK, False),
          ("próximos passos", BLUE, False),
          (" de vocês serão:", GRAY_DARK, False),
        ]
    bullets: list of str — itens da lista (Red Hat Display, #3C3C3C, 8.6pt).
      Exemplo:
        [
          "Baixar as ferramentas necessárias para o trabalho;",
          "Configuração do ambiente de desenvolvimento;",
          "Conseguir acesso ao Teams e Azure DevOps;",
        ]
    section_name: str — rodapé de seção.
    """
    slide = _clone_slide(prs, TMPL_BULLETS)

    # Texto introdutório (posição original: top ≈ 1.899in)
    intro_box = _find_textbox_near(slide, 0.833, 1.899)
    if intro_box:
        _set_spans(intro_box, intro_spans, size_pt=21.4)

    # Lista de bullets (posição original: top ≈ 2.964in)
    bullets_box = _find_textbox_near(slide, 0.833, 2.964)
    if bullets_box:
        _set_multiline(bullets_box, bullets, FONT_DISPLAY, 8.6, False, GRAY_DARK)

    # Rodapé (posição original: top ≈ 4.719in)
    footer = _find_textbox_near(slide, 0.833, 4.719)
    if footer:
        _set_simple(footer, section_name, FONT_MONO, 5.0, True, GRAY_DARK)

    return slide


# =============================================================================
# === CONTEÚDO DA APRESENTAÇÃO ================================================
# === Substitua apenas este bloco com o conteúdo real. =======================
# =============================================================================

prs = Presentation(TEMPLATE_PATH)
ORIGINAL_SLIDE_COUNT = len(prs.slides)  # 13 slides do template

# --- Slide 1: Citação de impacto ---
quote_slide(
    prs,
    spans=[
        ('"[Frase de impacto da apresentação com ', GRAY_DARK, False),
        ('palavra em destaque', BLUE, True),
        ('.]"', GRAY_DARK, False),
    ],
    section_name="[Nome da Seção]"
)

# --- Slide 2: Título + Corpo ---
content_slide(
    prs,
    title_line1="[Primeira linha",
    title_line2="do título?]",
    body_text="[Texto descritivo do slide. Use Inter para corpo de texto. "
              "Máximo de 4–5 linhas para manter a legibilidade da apresentação.]",
    section_name="[Nome da Seção]"
)

# --- Slide 3: Palavras-chave em azul ---
keywords_slide(
    prs,
    page_title="[Título da página]",
    subtitle_label="[Label do contexto:]",
    keywords=[
        "[Palavra 1]",
        "[Palavra 2]",
        "[Palavra 3]",
        "[Palavra 4]",
        "[Palavra 5]",
    ],
    body_text="[Texto de apoio em Inter explicando o contexto das palavras-chave.]",
    section_name="[Nome da Seção]"
)

# --- Slide 4: Bullets / Lista de ações ---
bullets_slide(
    prs,
    intro_spans=[
        ("[Texto introdutório com ", GRAY_DARK, False),
        ("destaque em azul", BLUE, False),
        (" no trecho relevante:]", GRAY_DARK, False),
    ],
    bullets=[
        "[Item da lista 1;]",
        "[Item da lista 2;]",
        "[Item da lista 3;]",
        "[Item da lista 4;]",
    ],
    section_name="[Nome da Seção]"
)

# --- Remover os slides originais do template (manter apenas os novos) ---
for _ in range(ORIGINAL_SLIDE_COUNT):
    _delete_slide(prs, 0)

# =============================================================================

prs.save(OUTPUT_PATH)
print(f"Apresentação gerada com sucesso: {OUTPUT_PATH}")
